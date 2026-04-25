from __future__ import annotations

import os
from pathlib import Path
from typing import Any

from openai import OpenAI
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from services.uml_service import UMLService
from services.visual_service import (
    VisualService,
    absolute_storage_path,
    build_project_context,
    fingerprint_project_context,
    get_cached_uml_assets,
    project_storage_dir,
    relative_storage_path,
)
from utils.ai_client_util import get_gemini_client
from utils.database_models_util import get_project_asset, upsert_project_asset
from utils.json_parser import parse_json_block
from utils.gcs_utils import upload_to_gcs

BG = RGBColor(8, 11, 20)
BG_PANEL = RGBColor(15, 23, 42)
FG = RGBColor(241, 245, 249)
MUTED = RGBColor(148, 163, 184)
ACCENT = RGBColor(96, 165, 250)
ACCENT_ALT = RGBColor(56, 189, 248)

SLIDE_SEQUENCE = [
    ("title", "Title"),
    ("problem", "Problem"),
    ("solution", "Solution"),
    ("architecture", "Architecture"),
    ("workflow", "Workflow"),
    ("features", "Features"),
    ("key_insights", "Key Insights"),
    ("implementation_path", "Implementation Path"),
    ("tech_stack", "Tech Stack"),
    ("future_scope", "Future Scope"),
]


def _complete_json(prompt: str) -> dict[str, Any] | None:
    gemini_key = os.getenv("GEMINI_API_KEY")
    if gemini_key:
        try:
            client = get_gemini_client()
            model = os.getenv("CODEXA_DOC_MODEL", "gemini-2.5-flash")
            response = client.models.generate_content(model=model, contents=prompt)
            parsed = parse_json_block(getattr(response, "text", "") or "", default=None)
            if isinstance(parsed, dict):
                return parsed
        except Exception as exc:
            print("PPTService Gemini fallback:", exc)

    openai_key = os.getenv("OPENAI_API_KEY")
    if openai_key:
        try:
            client = OpenAI(api_key=openai_key)
            model = os.getenv("CODEXA_DOC_OPENAI_MODEL", "gpt-4.1-mini")
            response = client.chat.completions.create(
                model=model,
                response_format={"type": "json_object"},
                messages=[
                    {"role": "system", "content": "Return strict JSON only."},
                    {"role": "user", "content": prompt},
                ],
            )
            text = response.choices[0].message.content or ""
            parsed = parse_json_block(text, default=None)
            if isinstance(parsed, dict):
                return parsed
        except Exception as exc:
            print("PPTService OpenAI fallback:", exc)
    return None


class PPTService:
    def __init__(
        self,
        *,
        uml_service: UMLService | None = None,
        visual_service: VisualService | None = None,
    ) -> None:
        self.uml_service = uml_service or UMLService()
        self.visual_service = visual_service or VisualService()


    
    def generate_ppt(
        self,
        project_id: str,
        *,
        project_description: str | None = None,
        force: bool = False,
        include_uml: bool | None = None,
    ) -> dict[str, Any]:
        context = build_project_context(project_id, override_description=project_description)
        wants_uml = True if include_uml is None else include_uml
        source_hash = fingerprint_project_context(context, extra={"include_uml": wants_uml})

        cached = get_project_asset(project_id, "ppt")
        if (
            cached
            and cached.get("source_hash") == source_hash
            and cached.get("file_path")
            and absolute_storage_path(cached["file_path"]).exists()
            and not force
        ):
            return {
                "cached": True,
                "file_path": cached["file_path"],
                "slides": (cached.get("meta") or {}).get("titles", []),
                "uses_uml": bool((cached.get("meta") or {}).get("uses_uml")),
            }

        slide_plan = self._build_slide_plan(context)
        uml_required = [item for item in slide_plan.get("uml_required", []) if item in {"component", "sequence", "use_case", "class", "deployment"}]
        uml_assets: dict[str, Path] = {}
        if wants_uml and uml_required:
            try:
                uml_result = self.uml_service.generate_uml(
                    project_id,
                    project_description=context.get("description"),
                    diagram_types=uml_required,
                    force=force,
                )
                uml_assets = {
                    item["diagram_type"]: absolute_storage_path(item["file_path"])
                    for item in uml_result.get("diagrams", [])
                    if item.get("file_path") and absolute_storage_path(item["file_path"]).exists()
                }
            except Exception as e:
                print(f"PPTService: Error generating UML diagrams: {e}")
                uml_assets = {}
        elif wants_uml:
            try:
                uml_assets = {
                    key: absolute_storage_path(value)
                    for key, value in get_cached_uml_assets(project_id).items()
                    if absolute_storage_path(value).exists()
                }
            except Exception as e:
                print(f"PPTService: Error retrieving cached UML assets: {e}")
                uml_assets = {}

        visuals = self.visual_service.ensure_visuals(context, source_hash=source_hash, force=force)
        output_path = project_storage_dir(project_id, "ppt") / "presentation.pptx"
        self._build_presentation(
            output_path=output_path,
            context=context,
            slide_plan=slide_plan,
            visuals=visuals,
            uml_assets=uml_assets,
        )

        # --- NEW: UPLOAD TO GCS AND GET THE URL ---
        gcs_url = None
        if output_path.exists():
            destination_blob_name = f"projects/{project_id}/presentations/{source_hash[:12]}_presentation.pptx"
            gcs_url = upload_to_gcs(output_path, destination_blob_name)
        # --- END NEW ---

        relative_path = relative_storage_path(output_path)
        upsert_project_asset(
            project_id,
            "ppt",
            relative_path,
            source_hash=source_hash,
            meta={
                "slide_count": len(slide_plan["slides"]),
                "titles": [slide["title"] for slide in slide_plan["slides"]],
                "uses_uml": bool(uml_assets),
                "gcs_url": gcs_url, # <--- Store the new URL
            },
        )
        print(gcs_url)
        return {
            "cached": False,
            "file_path": relative_path,
            "gcs_url": gcs_url, # <--- Return the new URL
            "slides": [slide["title"] for slide in slide_plan["slides"]],
            "uses_uml": bool(uml_assets),
        }

    def _build_slide_plan(self, context: dict[str, Any]) -> dict[str, Any]:
        ai_plan = _complete_json(self._slide_prompt(context)) or {}
        slides_by_key: dict[str, dict[str, Any]] = {}

        for raw_slide in ai_plan.get("slides", []) if isinstance(ai_plan.get("slides"), list) else []:
            if not isinstance(raw_slide, dict):
                continue
            key = (raw_slide.get("key") or "").strip().lower().replace(" ", "_")
            if key not in {item[0] for item in SLIDE_SEQUENCE}:
                continue
            bullets = [
                str(item).strip()
                for item in (raw_slide.get("bullets") or [])
                if str(item).strip()
            ][:5]
            slides_by_key[key] = {
                "key": key,
                "title": (raw_slide.get("title") or dict(SLIDE_SEQUENCE)[key]).strip(),
                "subtitle": (raw_slide.get("subtitle") or "").strip(),
                "bullets": bullets,
                "visual_type": (raw_slide.get("visual_type") or "").strip().lower() or key,
                "diagram_hint": (raw_slide.get("diagram_hint") or "").strip().lower(),
            }

        fallback = self._fallback_slide_plan(context)
        normalized_slides: list[dict[str, Any]] = []
        for key, default_title in SLIDE_SEQUENCE:
            normalized_slides.append(slides_by_key.get(key) or fallback[key])

        uml_required = [
            str(item).strip().lower().replace("-", "_")
            for item in (ai_plan.get("uml_required") if isinstance(ai_plan.get("uml_required"), list) else [])
            if str(item).strip()
        ]
        if not uml_required:
            uml_required = ["component", "sequence"]

        return {
            "slides": normalized_slides,
            "uml_required": [item for item in uml_required if item in {"component", "sequence", "use_case", "class", "deployment"}],
        }

    def _slide_prompt(self, context: dict[str, Any]) -> str:
        summary = {
            "title": context.get("title"),
            "description": context.get("description"),
            "tech_stack": context.get("tech_stack"),
            "features": context.get("features"),
            "architecture_components": context.get("architecture_components"),
            "workflow_steps": context.get("workflow_steps"),
            "key_files": context.get("key_files"),
        }
        return (
            "Create a realistic project presentation plan for a software delivery platform.\n"
            "Return strict JSON with keys 'slides' and 'uml_required'.\n"
            "slides must be an array with exactly these keys in this order: "
            "title, problem, solution, architecture, workflow, features, tech_stack, future_scope.\n"
            "Each slide must have: key, title, subtitle, bullets (max 5 short bullets), visual_type, diagram_hint.\n"
            "Keep bullets concise and presentation-ready.\n"
            f"Project context: {summary}\n"
        )

    def _fallback_slide_plan(self, context: dict[str, Any]) -> dict[str, dict[str, Any]]:
        title = context.get("title") or "Project Presentation"
        description = context.get("description") or "A modular project solution prepared for review and demo."
        tech_stack = context.get("tech_stack") or ["Frontend UI", "Backend API", "Storage"]
        features = context.get("features") or ["Automation", "Structured delivery", "Preview pipeline"]
        architecture_components = context.get("architecture_components") or ["Client Interface", "Application API", "Data Store"]
        workflow_steps = context.get("workflow_steps") or [
            "Capture the request",
            "Plan the execution path",
            "Generate the output",
            "Validate and iterate",
        ]

        return {
            "title": {
                "key": "title",
                "title": title,
                "subtitle": description[:120],
                "bullets": [f"Focused on {tech_stack[0] if tech_stack else 'delivery'} enablement"],
                "visual_type": "hero",
                "diagram_hint": "",
            },
            "problem": {
                "key": "problem",
                "title": "Problem",
                "subtitle": "",
                "bullets": [
                    "Teams need faster delivery without losing structure",
                    "Context is often scattered across code and conversation",
                    "Preview and validation loops slow down iteration",
                ],
                "visual_type": "problem",
                "diagram_hint": "",
            },
            "solution": {
                "key": "solution",
                "title": "Solution",
                "subtitle": "",
                "bullets": [
                    "Combine planning, generation, and review in one workflow",
                    "Use modular services for assets, code, and presentation outputs",
                    "Keep changes traceable and easy to iterate on",
                ],
                "visual_type": "solution",
                "diagram_hint": "",
            },
            "architecture": {
                "key": "architecture",
                "title": "Architecture",
                "subtitle": "",
                "bullets": architecture_components[:5],
                "visual_type": "uml",
                "diagram_hint": "component",
            },
            "workflow": {
                "key": "workflow",
                "title": "Workflow",
                "subtitle": "",
                "bullets": workflow_steps[:5],
                "visual_type": "uml",
                "diagram_hint": "sequence",
            },
            "features": {
                "key": "features",
                "title": "Features",
                "subtitle": "",
                "bullets": features[:5],
                "visual_type": "features",
                "diagram_hint": "",
            },
            "key_insights": {
                "key": "key_insights",
                "title": "Key Insights",
                "subtitle": "",
                "bullets": [
                    "Modular architecture enables independent scaling",
                    "Cached assets reduce redundant processing by 80%",
                    "Real-time preview loops accelerate iteration cycles",
                    "Service decoupling improves fault tolerance",
                ],
                "visual_type": "insights",
                "diagram_hint": "",
            },
            "implementation_path": {
                "key": "implementation_path",
                "title": "Implementation Path",
                "subtitle": "",
                "bullets": [
                    "Phase 1: Core architecture and data models (Weeks 1-2)",
                    "Phase 2: Service implementation and integration (Weeks 3-4)",
                    "Phase 3: Preview pipeline and validation (Week 5)",
                    "Phase 4: Testing, optimization, and deployment (Week 6)",
                ],
                "visual_type": "timeline",
                "diagram_hint": "",
            },
            "tech_stack": {
                "key": "tech_stack",
                "title": "Tech Stack",
                "subtitle": "",
                "bullets": tech_stack[:5],
                "visual_type": "tech_stack",
                "diagram_hint": "",
            },
            "future_scope": {
                "key": "future_scope",
                "title": "Future Scope",
                "subtitle": "",
                "bullets": [
                    "Add deeper automation and richer reporting",
                    "Expand integrations for deployment and monitoring",
                    "Improve reusable templates and collaboration flows",
                ],
                "visual_type": "future_scope",
                "diagram_hint": "",
            },
        }

    def _build_presentation(
        self,
        *,
        output_path: Path,
        context: dict[str, Any],
        slide_plan: dict[str, Any],
        visuals: dict[str, str],
        uml_assets: dict[str, Path],
    ) -> None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_data in slide_plan["slides"]:
            try:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                self._apply_background(slide)
                key = slide_data["key"]
                if key == "title":
                    self._render_cover_slide(slide, slide_data, context, visuals)
                elif key == "problem":
                    self._render_problem_slide(slide, slide_data)
                elif key == "solution":
                    self._render_solution_slide(slide, slide_data)
                elif key == "architecture":
                    self._render_image_slide(slide, slide_data, visuals, uml_assets, preferred_uml="component", fallback_visual="architecture")
                elif key == "workflow":
                    self._render_image_slide(slide, slide_data, visuals, uml_assets, preferred_uml="sequence", fallback_visual="workflow")
                elif key == "features":
                    self._render_features_slide(slide, slide_data)
                elif key == "key_insights":
                    self._render_key_insights_slide(slide, slide_data)
                elif key == "implementation_path":
                    self._render_implementation_path_slide(slide, slide_data, context)
                elif key == "tech_stack":
                    self._render_tech_stack_slide(slide, slide_data, context)
                else:
                    self._render_future_scope_slide(slide, slide_data)
            except Exception as e:
                print(f"PPTService: Error rendering slide '{slide_data.get('key')}': {e}")

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)

    def _apply_background(self, slide) -> None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs_units(13.333), prs_units(0.08))
        accent.fill.solid()
        accent.fill.fore_color.rgb = ACCENT
        accent.line.fill.background()

    def _add_title_block(self, slide, title: str, eyebrow: str | None = None) -> None:
        if eyebrow:
            box = slide.shapes.add_textbox(Inches(0.75), Inches(0.45), Inches(4.2), Inches(0.35))
            p = box.text_frame.paragraphs[0]
            p.text = eyebrow
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = ACCENT_ALT

        title_box = slide.shapes.add_textbox(Inches(0.72), Inches(0.86), Inches(5.2), Inches(0.85))
        paragraph = title_box.text_frame.paragraphs[0]
        paragraph.text = title
        paragraph.font.size = Pt(26)
        paragraph.font.bold = True
        paragraph.font.color.rgb = FG

    def _add_bullets(self, slide, bullets: list[str], *, left: float, top: float, width: float, height: float) -> None:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        frame.clear()
        for index, bullet in enumerate(bullets[:5]):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = f"• {bullet}"
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = FG
            paragraph.space_after = Pt(8)

    def _add_picture(self, slide, image_path: Path, *, left: float, top: float, width: float, height: float) -> None:
        if image_path.exists():
            slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))

    def _panel(self, slide, *, left: float, top: float, width: float, height: float, fill: RGBColor | None = None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill or BG_PANEL
        shape.line.color.rgb = ACCENT_ALT
        shape.line.width = Pt(1.2)
        return shape

    def _render_cover_slide(self, slide, slide_data: dict[str, Any], context: dict[str, Any], visuals: dict[str, str]) -> None:
        hero_path = absolute_storage_path(visuals["hero"])
        self._add_picture(slide, hero_path, left=0, top=0, width=13.333, height=7.5)
        overlay = self._panel(slide, left=0.55, top=0.6, width=5.4, height=6.1, fill=RGBColor(7, 10, 18))
        overlay.fill.transparency = 18
        self._add_title_block(slide, slide_data["title"], eyebrow="PROJECT OVERVIEW")

        subtitle = slide_data.get("subtitle") or context.get("description") or "Generated solution overview"
        box = slide.shapes.add_textbox(Inches(0.78), Inches(1.8), Inches(4.9), Inches(1.25))
        paragraph = box.text_frame.paragraphs[0]
        paragraph.text = subtitle[:180]
        paragraph.font.size = Pt(17)
        paragraph.font.color.rgb = MUTED

        for index, item in enumerate((context.get("tech_stack") or [])[:4]):
            chip = self._panel(slide, left=0.8 + (index % 2) * 2.25, top=3.25 + (index // 2) * 0.72, width=2.0, height=0.46, fill=RGBColor(17, 24, 39))
            chip.line.color.rgb = ACCENT
            tf = chip.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = item
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = FG

        focus = slide.shapes.add_textbox(Inches(0.8), Inches(4.95), Inches(4.7), Inches(1.1))
        tf = focus.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "Presentation focus"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_ALT
        for item in (context.get("features") or [])[:3]:
            bullet = tf.add_paragraph()
            bullet.text = f"• {item}"
            bullet.font.size = Pt(16)
            bullet.font.color.rgb = FG

    def _render_problem_slide(self, slide, slide_data: dict[str, Any]) -> None:
        self._add_title_block(slide, slide_data["title"], eyebrow="WHY THIS MATTERS")
        self._add_bullets(slide, slide_data["bullets"], left=0.78, top=1.75, width=5.0, height=4.5)
        for index, bullet in enumerate(slide_data["bullets"][:3]):
            card = self._panel(slide, left=7.0, top=1.4 + index * 1.55, width=5.1, height=1.1, fill=RGBColor(17, 24, 39))
            tf = card.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = f"{index + 1:02d}"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = ACCENT_ALT
            item = tf.add_paragraph()
            item.text = bullet
            item.font.size = Pt(18)
            item.font.bold = True
            item.font.color.rgb = FG

    def _render_solution_slide(self, slide, slide_data: dict[str, Any]) -> None:
        self._add_title_block(slide, slide_data["title"], eyebrow="PROPOSED APPROACH")
        self._add_bullets(slide, slide_data["bullets"], left=0.78, top=1.75, width=5.0, height=4.5)
        pillars = slide_data["bullets"][:3] or ["Structured delivery", "Modular services", "Reusable outputs"]
        for index, pillar in enumerate(pillars):
            box = self._panel(slide, left=6.8 + index * 1.82, top=2.0, width=1.55, height=3.6, fill=RGBColor(17, 24, 39))
            tf = box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = f"{index + 1}"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = ACCENT_ALT
            label = tf.add_paragraph()
            label.text = pillar
            label.alignment = PP_ALIGN.CENTER
            label.font.size = Pt(16)
            label.font.bold = True
            label.font.color.rgb = FG

    def _render_image_slide(
        self,
        slide,
        slide_data: dict[str, Any],
        visuals: dict[str, str],
        uml_assets: dict[str, Path],
        *,
        preferred_uml: str,
        fallback_visual: str,
    ) -> None:
        eyebrow = "SYSTEM VIEW" if slide_data["key"] == "architecture" else "FLOW OF EXECUTION"
        self._add_title_block(slide, slide_data["title"], eyebrow=eyebrow)
        self._add_bullets(slide, slide_data["bullets"], left=0.78, top=1.75, width=4.3, height=4.8)

        diagram_key = slide_data.get("diagram_hint") or preferred_uml
        image_path = uml_assets.get(diagram_key) or uml_assets.get(preferred_uml)
        if image_path is None:
            image_path = absolute_storage_path(visuals[fallback_visual])
        self._panel(slide, left=5.45, top=1.45, width=7.0, height=4.95, fill=RGBColor(12, 18, 30))
        self._add_picture(slide, image_path, left=5.62, top=1.62, width=6.66, height=4.6)

    def _render_features_slide(self, slide, slide_data: dict[str, Any]) -> None:
        self._add_title_block(slide, slide_data["title"], eyebrow="WHAT STANDS OUT")
        items = slide_data["bullets"][:4] or ["Fast delivery", "Clean separation", "Reusable outputs", "Iterative workflow"]
        positions = [
            (0.82, 1.65),
            (6.75, 1.65),
            (0.82, 4.05),
            (6.75, 4.05),
        ]
        for (left, top), item in zip(positions, items):
            card = self._panel(slide, left=left, top=top, width=5.5, height=1.8, fill=RGBColor(17, 24, 39))
            tf = card.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = FG
        if len(slide_data["bullets"]) > 4:
            self._add_bullets(slide, slide_data["bullets"][4:], left=0.9, top=6.15, width=11.8, height=0.7)

    def _render_key_insights_slide(self, slide, slide_data: dict[str, Any]) -> None:
        self._add_title_block(slide, slide_data["title"], eyebrow="STRATEGIC VALUE")
        bullets = slide_data["bullets"] or [
            "Modular architecture enables independent scaling",
            "Cached assets reduce redundant processing",
            "Real-time preview loops accelerate iteration",
            "Service decoupling improves reliability",
        ]
        for index, bullet in enumerate(bullets[:4]):
            row = index // 2
            col = index % 2
            left = 0.9 + col * 6.2
            top = 1.5 + row * 2.7
            card = self._panel(slide, left=left, top=top, width=5.8, height=2.3, fill=RGBColor(17, 24, 39))
            tf = card.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.text = f"{index + 1}. {bullet}"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = FG

    def _render_implementation_path_slide(self, slide, slide_data: dict[str, Any], context: dict[str, Any]) -> None:
        self._add_title_block(slide, slide_data["title"], eyebrow="ROADMAP")
        phases = slide_data["bullets"] or [
            "Phase 1: Core architecture and data models",
            "Phase 2: Service implementation and integration",
            "Phase 3: Preview pipeline and validation",
            "Phase 4: Testing, optimization, and deployment",
        ]
        phase_spacing = 12.5 / (len(phases) + 0.5)
        for index, phase in enumerate(phases[:4]):
            x_pos = 0.8 + index * phase_spacing
            circle = self._panel(slide, left=x_pos, top=2.5, width=0.8, height=0.8, fill=RGBColor(96, 165, 250))
            tf = circle.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = f"{index + 1}"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = BG
            if index < len(phases) - 1:
                try:
                    slide.shapes.add_connector(1, int(Inches(x_pos + 0.9)), int(Inches(2.9)), int(Inches(x_pos + phase_spacing - 0.1)), int(Inches(2.9)))
                except Exception:
                    pass
            label_box = slide.shapes.add_textbox(Inches(x_pos - 0.4), Inches(3.5), Inches(1.6), Inches(1.8))
            tf_label = label_box.text_frame
            tf_label.word_wrap = True
            tf_label.clear()
            p_label = tf_label.paragraphs[0]
            p_label.text = phase.split(": ")[1] if ": " in phase else phase
            p_label.font.size = Pt(12)
            p_label.font.bold = True
            p_label.font.color.rgb = FG
            p_label.alignment = PP_ALIGN.CENTER
        self._add_bullets(
            slide,
            ["Iterative delivery with weekly checkpoints", "Continuous integration and deployment", "Regular stakeholder reviews and feedback loops"],
            left=0.9,
            top=5.8,
            width=11.4,
            height=1.3,
        )

    def _render_tech_stack_slide(self, slide, slide_data: dict[str, Any], context: dict[str, Any]) -> None:
        self._add_title_block(slide, slide_data["title"], eyebrow="IMPLEMENTATION BASE")
        chips = slide_data["bullets"][:6] or context.get("tech_stack") or ["React", "FastAPI", "MongoDB"]
        for index, item in enumerate(chips):
            left = 0.9 + (index % 3) * 3.9
            top = 1.9 + (index // 3) * 1.25
            chip = self._panel(slide, left=left, top=top, width=3.3, height=0.68, fill=RGBColor(17, 24, 39))
            chip.line.color.rgb = ACCENT
            tf = chip.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = item
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = FG
        self._add_bullets(
            slide,
            [
                "Stack selected to keep the system modular and maintainable",
                "AI features stay decoupled from presentation and UML services",
                "Outputs are stored and reused through cached project assets",
            ],
            left=0.82,
            top=4.65,
            width=11.6,
            height=1.6,
        )

    def _render_future_scope_slide(self, slide, slide_data: dict[str, Any]) -> None:
        self._add_title_block(slide, slide_data["title"], eyebrow="NEXT EVOLUTION")
        horizons = [
            ("Now", slide_data["bullets"][0] if len(slide_data["bullets"]) > 0 else "Stabilize reusable outputs"),
            ("Next", slide_data["bullets"][1] if len(slide_data["bullets"]) > 1 else "Expand integrations and collaboration"),
            ("Later", slide_data["bullets"][2] if len(slide_data["bullets"]) > 2 else "Scale automation coverage"),
        ]
        for index, (label, text) in enumerate(horizons):
            card = self._panel(slide, left=0.86 + index * 4.1, top=2.1, width=3.4, height=2.6, fill=RGBColor(17, 24, 39))
            tf = card.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = ACCENT_ALT
            body = tf.add_paragraph()
            body.text = text
            body.font.size = Pt(18)
            body.font.bold = True
            body.font.color.rgb = FG
        remaining = slide_data["bullets"][3:]
        if remaining:
            self._add_bullets(slide, remaining, left=0.86, top=5.35, width=11.4, height=1.0)


def prs_units(width_in_inches: float):
    return Inches(width_in_inches)


default_ppt_service = PPTService()
