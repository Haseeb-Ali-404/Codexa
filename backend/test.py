import os
import requests
import zlib
import json
import re
from concurrent.futures import ThreadPoolExecutor
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from openai import OpenAI
from dotenv import load_dotenv

load_dotenv()

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

OUTPUT_DIR = "styled_output"
PLANTUML_SERVER = "https://www.plantuml.com/plantuml/png/"

os.makedirs(OUTPUT_DIR, exist_ok=True)
client = OpenAI(api_key=OPENAI_API_KEY)

# ---------------- STATIC BLACK & WHITE THEME ----------------
theme = {
    "mode": "dark",
    "primary": "#ffffff",
    "secondary": "#1a1a1a",
    "background": "#000000",
    "text": "#ffffff"
}

# ---------------- PROJECT DESCRIPTION ----------------
PROJECT_DESCRIPTION = """
SmartHire is a full-stack AI-powered recruitment and hiring platform designed to streamline the entire hiring lifecycle for companies.

The platform supports three main user roles: administrators, recruiters, and candidates. Administrators manage system-wide settings, monitor activity logs, and oversee platform performance. Recruiters are responsible for creating job postings, reviewing applications, shortlisting candidates, scheduling interviews, and tracking hiring progress. Candidates can create profiles, upload resumes, apply for jobs, and track the status of their applications.

The system includes a secure authentication module with role-based access control, ensuring that each user interacts only with the features relevant to their role. Authentication is implemented using token-based mechanisms, and user sessions are securely managed.

Recruiters can create job listings by specifying role details, required skills, experience level, and job descriptions. These listings are stored in the database and made available to candidates. Candidates can search and filter job listings based on categories, location, and skills, and then submit applications.

Once a candidate applies, the system creates an application record linking the candidate and the job. Recruiters can view applications, update their status (such as applied, shortlisted, interviewed, or rejected), and add internal notes.

An AI-based resume analysis module processes uploaded resumes and extracts key information such as skills, education, and experience. This data is used to automatically match candidates with suitable job postings and rank applications based on relevance.

The platform also includes an interview scheduling system where recruiters can schedule interviews, send notifications, and manage interview timelines. Candidates receive notifications and can confirm or reschedule interviews.

A notification system is implemented to handle events such as application status updates, interview schedules, and important alerts. Notifications are delivered via email and in-app messages.

The backend architecture follows a modular service-based design, with separate layers for API routing, business logic, and data access. The system uses well-defined schemas for users, jobs, applications, resumes, and notifications, ensuring efficient relationships and scalability.

The frontend is built using a modern component-based architecture, featuring dashboards for recruiters and candidates. It includes interactive UI elements such as tables, filters, and charts for analytics like hiring trends, application success rates, and recruiter performance.

The platform is designed to be scalable and extensible, allowing future integration of advanced AI features such as automated interview analysis, candidate scoring, and predictive hiring insights.
"""

# ---------------- SAFE JSON PARSER ----------------
def safe_json_load(text):
    text = text.strip()
    text = re.sub(r"```json|```", "", text)

    match = re.search(r'(\{.*\}|\[.*\])', text, re.DOTALL)
    if not match:
        raise Exception("❌ No JSON found:\n" + text)

    return json.loads(match.group(0))

# ---------------- UML STYLE ----------------
def build_uml_style():
    return """
skinparam backgroundColor #000000
skinparam defaultFontName Inter
skinparam defaultFontSize 14

skinparam ArrowColor #ffffff
skinparam ArrowThickness 2
skinparam shadowing false
skinparam roundcorner 12
skinparam linetype ortho
skinparam dpi 150

skinparam class {
  BackgroundColor #1a1a1a
  BorderColor #ffffff
  FontColor #ffffff
}

skinparam component {
  BackgroundColor #1a1a1a
  BorderColor #ffffff
  FontColor #ffffff
}

skinparam actor {
  BackgroundColor #1a1a1a
  BorderColor #ffffff
  FontColor #ffffff
}

skinparam usecase {
  BackgroundColor #1a1a1a
  BorderColor #ffffff
  FontColor #ffffff
}

skinparam package {
  BorderColor #ffffff
  FontColor #ffffff
}

scale 1.2
"""

# ---------------- APPLY STYLE ----------------
def apply_style(code):
    style = build_uml_style()
    return code.replace("@startuml", f"@startuml\n{style}")

# ---------------- ENCODE PLANTUML ----------------
def encode_plantuml(text):
    data = zlib.compress(text.encode("utf-8"))[2:-4]
    return encode_base64(data)

def encode_base64(data):
    alphabet = "0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz-_"
    result = ""

    def append3(b1, b2, b3):
        return (
            alphabet[b1 >> 2] +
            alphabet[((b1 & 0x3) << 4) | (b2 >> 4)] +
            alphabet[((b2 & 0xF) << 2) | (b3 >> 6)] +
            alphabet[b3 & 0x3F]
        )

    i = 0
    while i < len(data):
        if i + 2 == len(data):
            result += append3(data[i], data[i+1], 0)
            break
        elif i + 1 == len(data):
            result += append3(data[i], 0, 0)
            break
        else:
            result += append3(data[i], data[i+1], data[i+2])
        i += 3

    return result

# ---------------- GENERATE UML ----------------
def generate_diagrams():
    prompt = f"""
    Generate 6 UML diagrams in PlantUML:

    {PROJECT_DESCRIPTION}

    Return JSON:
    {{
      "use_case": "...",
      "class": "...",
      "sequence": "...",
      "component": "...",
      "deployment": "...",
      "activity": "..."
    }}

    Only JSON.
    """

    res = client.chat.completions.create(
        model="gpt-4o",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3
    )

    return safe_json_load(res.choices[0].message.content)

# ---------------- IMAGE GENERATION ----------------
def generate_image(name, code):
    styled_code = apply_style(code)
    encoded = encode_plantuml(styled_code)
    url = PLANTUML_SERVER + encoded

    path = os.path.join(OUTPUT_DIR, f"{name}.png")

    img = requests.get(url).content
    with open(path, "wb") as f:
        f.write(img)

    return path

# ---------------- PARALLEL IMAGE ----------------
def generate_all_images(diagrams):
    results = {}

    def worker(item):
        name, code = item
        return name, generate_image(name, code)

    with ThreadPoolExecutor(max_workers=6) as executor:
        for name, path in executor.map(worker, diagrams.items()):
            results[name] = path

    return results

# ---------------- PPT CONTENT ----------------
def generate_ppt_content():
    prompt = f"""
    Create a professional PPT for:

    {PROJECT_DESCRIPTION}

    Return JSON:
    [
      {{ "title": "...", "points": ["...", "..."] }}
    ]

    Include:
    Title, Problem, Solution, Features,
    Architecture, UML explanation, Tech stack, Conclusion

    Only JSON.
    """

    res = client.chat.completions.create(
        model="gpt-4o",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.5
    )

    return safe_json_load(res.choices[0].message.content)

# ---------------- CREATE PPT ----------------
def create_ppt(slides, images):
    prs = Presentation()

    bg_color = RGBColor(0, 0, 0)
    text_color = RGBColor(255, 255, 255)

    for i, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # Background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color

        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = slide_data["title"]
        content.text = "\n".join([f"• {p}" for p in slide_data["points"]])

        # Apply text color
        title.text_frame.paragraphs[0].font.color.rgb = text_color
        content.text_frame.paragraphs[0].font.color.rgb = text_color

        # Add diagram
        if i < len(images):
            slide.shapes.add_picture(
                list(images.values())[i],
                Inches(1),
                Inches(3),
                width=Inches(6)
            )

    ppt_path = os.path.join(OUTPUT_DIR, "black_white_presentation.pptx")
    prs.save(ppt_path)

    return ppt_path

# ---------------- MAIN ----------------
def main():
    print("🚀 Generating UML diagrams...")
    diagrams = generate_diagrams()

    print("⚡ Rendering diagrams...")
    images = generate_all_images(diagrams)

    print("🧠 Generating PPT content...")
    slides = generate_ppt_content()

    print("📊 Creating PPT...")
    ppt = create_ppt(slides, images)    

    print("\n✅ DONE")
    print("Output folder:", OUTPUT_DIR)
    print("PPT:", ppt)

if __name__ == "__main__":
    main()